import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

BASE_DIR = "/home/villarias/Projects/jayraldinescateringsystem"
PPTX_OUT = os.path.join(BASE_DIR, "PROJECT_PROPOSAL_PRESENTATION.pptx")

# Color Palette Constants
COLOR_BG_DARK = RGBColor(15, 23, 42)      # #0F172A Slate 900
COLOR_CARD_DARK = RGBColor(30, 41, 59)    # #1E293B Slate 800
COLOR_CRIMSON = RGBColor(225, 29, 72)     # #E11D48 Accent Rose/Crimson
COLOR_GOLD = RGBColor(245, 158, 11)       # #F59E0B Amber/Gold
COLOR_TEXT_LIGHT = RGBColor(248, 250, 252)# #F8FAFC
COLOR_TEXT_MUTED = RGBColor(148, 163, 184)# #94A3B8
COLOR_CARD_BORDER = RGBColor(51, 65, 85)  # #334155
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_EMERALD = RGBColor(16, 185, 129)    # #10B981

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    def add_base_slide(title_text, subtitle_text=None, presenter_badge=None):
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG_DARK
        bg.line.fill.background()

        # Top Accent Bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_CRIMSON
        bar.line.fill.background()

        # Title & Subtitle Box
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = "Arial"
            p2.font.size = Pt(12)
            p2.font.color.rgb = COLOR_TEXT_MUTED
            p2.space_before = Pt(3)

        # Presenter Badge if present
        if presenter_badge:
            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.2), Inches(0.45), Inches(2.3), Inches(0.45))
            badge.fill.solid()
            badge.fill.fore_color.rgb = COLOR_CARD_DARK
            badge.line.color.rgb = COLOR_CRIMSON
            badge.line.width = Pt(1.5)
            btf = badge.text_frame
            bp = btf.paragraphs[0]
            bp.text = presenter_badge
            bp.font.name = "Arial"
            bp.font.size = Pt(10)
            bp.font.bold = True
            bp.font.color.rgb = COLOR_GOLD
            bp.alignment = PP_ALIGN.CENTER

        # Footer
        footer = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.3))
        ftf = footer.text_frame
        ftf.margin_left = ftf.margin_top = ftf.margin_right = ftf.margin_bottom = 0
        fp = ftf.paragraphs[0]
        fp.text = "Jayraldine's Catering Management & Self-Service Kiosk System  |  Formal Project Proposal"
        fp.font.name = "Arial"
        fp.font.size = Pt(9)
        fp.font.color.rgb = RGBColor(100, 116, 139)

        return slide

    def add_card(slide, left, top, width, height, title, items, badge=None, accent_color=COLOR_CRIMSON):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_DARK
        card.line.color.rgb = COLOR_CARD_BORDER
        card.line.width = Pt(1)

        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = accent_color

        for item in items:
            p = tf.add_paragraph()
            p.text = f"•  {item}"
            p.font.name = "Arial"
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_TEXT_LIGHT
            p.space_before = Pt(6)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 1: Title Slide
    # ─────────────────────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank_slide_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_BG_DARK
    bg1.line.fill.background()

    # Center hero box
    tb1 = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.933), Inches(4.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0

    p = tf1.paragraphs[0]
    p.text = "JAYRALDINE'S CATERING MANAGEMENT\n& SELF-SERVICE KIOSK SYSTEM"
    p.font.name = "Arial"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_CRIMSON
    p.alignment = PP_ALIGN.CENTER

    p2 = tf1.add_paragraph()
    p2.text = "Formal Technical Project Proposal & System Specification"
    p2.font.name = "Arial"
    p2.font.size = Pt(16)
    p2.font.color.rgb = COLOR_GOLD
    p2.space_before = Pt(14)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf1.add_paragraph()
    p3.text = "Client: Jayraldine's Catering Services (Cebu, Philippines)\nPresented by: Project Development Team (3 Presenters)"
    p3.font.name = "Arial"
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_TEXT_MUTED
    p3.space_before = Pt(18)
    p3.alignment = PP_ALIGN.CENTER

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 2: Executive Summary & Project Background (Presenter 1)
    # ─────────────────────────────────────────────────────────────────────────
    s2 = add_base_slide("Executive Summary & Project Scope", "High-level overview of the dual-platform catering ecosystem", "PRESENTER 1")
    add_card(s2, Inches(0.8), Inches(1.7), Inches(5.6), Inches(4.9), "Client & Objective", [
        "Client: Jayraldine's Catering Services (Cebu, PH)",
        "Classification: Enterprise Catering ERP + Touch Kiosk",
        "Primary Goal: Automate manual paperwork, order estimation, and financial accounting.",
        "Target Users: Business Owners, Event Coordinators, Kitchen Staff, and Catering Clients."
    ], accent_color=COLOR_GOLD)

    add_card(s2, Inches(6.9), Inches(1.7), Inches(5.6), Inches(4.9), "Dual-Platform Architecture", [
        "Windows PC Suite (Catering_Present): Back-office control center for billing, invoices, calendar, expenses, and AI.",
        "Portable Field Kiosk (Tablet / APK): 100% offline self-service ordering terminal for bridal fairs, tastings, and venues.",
        "2-Way Data Synchronization: Seamless master sync via Excel templates and SQLite database export/import.",
        "Zero Subscription Fees: Local data ownership with zero monthly cloud SaaS costs."
    ], accent_color=COLOR_CRIMSON)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 3: Problems vs Solutions (Presenter 1)
    # ─────────────────────────────────────────────────────────────────────────
    s3 = add_base_slide("Operational Problems vs Proposed Solutions", "Direct business impact and problem resolution matrix", "PRESENTER 1")
    add_card(s3, Inches(0.8), Inches(1.7), Inches(3.6), Inches(4.9), "Manual Inefficiencies", [
        "Math & Headcount Errors: Handwritten estimation causes balance and pricing mistakes.",
        "Unmonitored Expenses: Revenue tracked without Cost of Goods Sold (COGS).",
        "Overbooking Risks: Schedule conflicts and excess capacity overload.",
        "Remote Disconnectivity: Cloud-only apps fail in remote resort venues."
    ], accent_color=RGBColor(239, 68, 68))

    add_card(s3, Inches(4.8), Inches(1.7), Inches(3.6), Inches(4.9), "System Solutions", [
        "Automated Pricing Engine: Exact mathematical calculations for packages, pax, and add-ons.",
        "Net Profit Tracking: Direct deduction of food, labor, and logistics from revenue.",
        "Capacity Hard-Block: Enforces max daily pax limits (e.g. 600 pax/day).",
        "100% Offline-First: Embedded SQLite WASM requires zero internet."
    ], accent_color=COLOR_EMERALD)

    add_card(s3, Inches(8.8), Inches(1.7), Inches(3.7), Inches(4.9), "Business Value & ROI", [
        "Speed: Reduces client booking consultation from 45 mins to <5 mins.",
        "Accuracy: 100% error-free invoicing, balance records, and down payments.",
        "Brand Elevation: Sleek touch kiosk builds prestige and client trust.",
        "Full Data Ownership: 100% private SQLite database on physical hardware."
    ], accent_color=COLOR_GOLD)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 4: Executive Dashboard & CRM Loyalty (Presenter 1)
    # ─────────────────────────────────────────────────────────────────────────
    s4 = add_base_slide("Executive Dashboard & Customer CRM", "Real-time oversight, booking pipeline, and loyalty tracking", "PRESENTER 1")
    add_card(s4, Inches(0.8), Inches(1.7), Inches(5.6), Inches(4.9), "Executive Dashboard KPIs", [
        "Real-Time Revenue: Live gross sales and payment collection ledgers.",
        "Net Profit Margin: Automated computation against operating expenses.",
        "Pipeline Tracker: Visual counts for Inquiries, Confirmed, In-Prep, and Completed.",
        "Capacity Load Alert: Real-time monitoring of daily guest headcount capacity."
    ], accent_color=COLOR_GOLD)

    add_card(s4, Inches(6.9), Inches(1.7), Inches(5.6), Inches(4.9), "Customer CRM & Loyalty Engine", [
        "Automated Loyalty Tiers: Customer upgrade matrix (Bronze: 1-2, Silver: 3-5, Gold: 6-9, VIP: 10+ events).",
        "Follow-up Reminder System: Scheduled callbacks and repeat event notifications.",
        "Client History: Comprehensive log of past bookings, venues, and custom dish preferences.",
        "Directory Search: Instant lookup of existing customers during ordering."
    ], accent_color=COLOR_CRIMSON)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 5: 6-Step Guided Kiosk Ordering (Presenter 2)
    # ─────────────────────────────────────────────────────────────────────────
    s5 = add_base_slide("Client Ordering Subsystem (6-Step Wizard)", "Touch-first guided ordering flow on Tablet and Mobile Kiosks", "PRESENTER 2")
    add_card(s5, Inches(0.8), Inches(1.7), Inches(3.6), Inches(4.9), "Steps 1 & 2: Setup & Package", [
        "Step 1 (Customer & Location): Searchable Cebu Municipality/Barangay selector with separated street field.",
        "Step 2 (Event & Package): Touch occasion picker (Weddings, Birthdays, Corporate).",
        "Package Inclusions Viewer: Displays buffet rates, minimum pax, and equipment (chafing dishes, skirting, servers)."
    ], accent_color=COLOR_GOLD)

    add_card(s5, Inches(4.8), Inches(1.7), Inches(3.6), Inches(4.9), "Steps 3 & 4: Menu & Add-ons", [
        "Step 3 (Categorized Menu): Grouped dishes (Beef, Pork, Chicken, Seafood, Desserts, Drinks).",
        "Sticky Navigation: Pinned category filter bar stays fixed while scrolling.",
        "Step 4 (Add-ons & Equipment): One-tap upsells for Lechon platters, dessert bars, and sound systems."
    ], accent_color=COLOR_CRIMSON)

    add_card(s5, Inches(8.8), Inches(1.7), Inches(3.7), Inches(4.9), "Steps 5 & 6: Terms & Receipt", [
        "Step 5 (Billing & Terms): Downpayment tracker (Cash, GCash, Bank) and catering terms agreement.",
        "Step 6 (Pre-Booking Verification): Summary confirmation drawer before committing.",
        "Instant Offline PDF Receipt: Itemized receipt generated and downloaded directly to device storage."
    ], accent_color=COLOR_EMERALD)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 6: Billing, Invoicing & Notifications (Presenter 2)
    # ─────────────────────────────────────────────────────────────────────────
    s6 = add_base_slide("Billing, Invoices & Notification Gateway", "Automated billing enforcement, payment tracking, and messaging", "PRESENTER 2")
    add_card(s6, Inches(0.8), Inches(1.7), Inches(5.6), Inches(4.9), "Commercial Invoicing & Downpayment", [
        "Official Invoicing: Automated invoice numbering (INV-XXXX) with itemized breakdown.",
        "Downpayment Hard-Lock: Blocks event confirmation without verified downpayment reception (e.g. 30%).",
        "Payment Milestones: Tracks initial deposits, partial balance collections, and settlement dates.",
        "Discounts & Surcharges: Configurable line-item adjustments and tax calculations."
    ], accent_color=COLOR_GOLD)

    add_card(s6, Inches(6.9), Inches(1.7), Inches(5.6), Inches(4.9), "Multi-Channel Notification Gateway", [
        "Automated SMTP Email Dispatcher: Sends booking confirmations and PDF invoices to customer email.",
        "SMS Gateway Integration: Real-time SMS alerts via Semaphore API upon approval.",
        "Receipt Printing: One-click high-resolution PDF export and commercial printing.",
        "Communication Logs: Audit history of all dispatched notifications per booking."
    ], accent_color=COLOR_CRIMSON)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 7: Cash Flow & Expense Management (Presenter 3)
    # ─────────────────────────────────────────────────────────────────────────
    s7 = add_base_slide("Cash Flow & Expense Tracking (Net Profit)", "Comprehensive financial control and true profitability accounting", "PRESENTER 3")
    add_card(s7, Inches(0.8), Inches(1.7), Inches(5.6), Inches(4.9), "Real-Time Cash Flow Monitor", [
        "Cash Inflow vs Outflow: Tracks liquid cash, GCash, and bank transfer receipts.",
        "Liquidity Balances: Real-time net cash-on-hand calculation.",
        "Payment Breakdown: Visual distribution of revenue across payment channels.",
        "Daily Financial Reconciliation: End-of-day transaction balancing."
    ], accent_color=COLOR_GOLD)

    add_card(s7, Inches(6.9), Inches(1.7), Inches(5.6), Inches(4.9), "Expense Management & COGS", [
        "Itemized Categories: Cost of Goods Sold (COGS), Food Ingredients, Labor/Waitstaff, Logistics/Fuel, Utilities, and Rentals.",
        "Net Profit Calculation: True automated formula: Gross Revenue - Total Expenses = Net Profit.",
        "Expense Logging: Date, category, description, and attached invoice amounts.",
        "Monthly Profit Reports: Comparative profit margins across fiscal periods."
    ], accent_color=COLOR_EMERALD)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 8: AI Intelligence & Demand Forecasting (Presenter 3)
    # ─────────────────────────────────────────────────────────────────────────
    s8 = add_base_slide("AI Intelligence & Smart Forecasting", "Machine learning insights for demand prediction and portion control", "PRESENTER 3")
    add_card(s8, Inches(0.8), Inches(1.7), Inches(5.6), Inches(4.9), "AI Demand & Volume Forecasting", [
        "Seasonal Prediction: Machine learning forecasting on peak wedding, holiday, and debut months.",
        "Expected Guest Turnout: Historical trend analysis predicting peak headcount surges.",
        "Capacity Planning: Recommends staffing and equipment allocation per weekend.",
        "Revenue Projections: Forecasts upcoming quarterly revenue based on booking velocity."
    ], accent_color=COLOR_CRIMSON)

    add_card(s8, Inches(6.9), Inches(1.7), Inches(5.6), Inches(4.9), "Smart Portion & Menu Analytics", [
        "Ingredient Estimator: Calculates raw food requirements based on pax count and selected dishes to prevent waste.",
        "Menu Engineering: Identifies high-margin signature dishes vs underperforming menu items.",
        "Cost Optimization: Recommends recipe portion adjustments to protect profit margins.",
        "Intelligent Insights: Actionable alerts for package pricing optimization."
    ], accent_color=COLOR_GOLD)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 9: Kitchen Operations & Kanban Workflow (Presenter 3)
    # ─────────────────────────────────────────────────────────────────────────
    s9 = add_base_slide("Kitchen Operations & Preparation Kanban", "Production pipeline management and quality assurance checklists", "PRESENTER 3")
    add_card(s9, Inches(0.8), Inches(1.7), Inches(5.6), Inches(4.9), "Kitchen Kanban Pipeline", [
        "Visual Order Pipeline: Real-time stages (Pending Prep, Cooking, Plating, Ready for Transport, Dispatched).",
        "Event Detail Cards: Displays event date, call time, venue location, and package tier.",
        "Status Transitions: Drag-and-drop / single-click order progression.",
        "Dispatch Coordination: Syncs kitchen readiness with logistical transport teams."
    ], accent_color=COLOR_GOLD)

    add_card(s9, Inches(6.9), Inches(1.7), Inches(5.6), Inches(4.9), "Per-Dish Preparation Checklist", [
        "Recipe Task Breakdown: Interactive sub-task checklists for every dish in the booking.",
        "Quality Standards: Ensures all appetizers, mains, and desserts meet standard execution.",
        "Zero Missed Inclusions: Prevents forgotten sauces, side dishes, or tableware.",
        "Staff Accountability: Records completion timestamps per kitchen station."
    ], accent_color=COLOR_CRIMSON)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 10: Complete 13-Module System Matrix
    # ─────────────────────────────────────────────────────────────────────────
    s10 = add_base_slide("Complete 13-Module System Functional Matrix", "Comprehensive ERP and Kiosk capability overview", "FULL SYSTEM")
    
    # Table layout
    rows, cols = 14, 3
    left, top, width, height = Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.1)
    table_shape = s10.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(6.6)
    table.columns[2].width = Inches(2.333)

    headers = ["Module", "Core Functionality", "Platform Availability"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_CRIMSON
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

    modules_data = [
        ("Executive Dashboard", "Real-time KPIs, Revenue, Net Profit, Capacity alerts, Pipeline", "PC Back-Office"),
        ("Booking & Kiosk Wizard", "6-step guided ordering, package selection, downpayment calculation", "PC + Tablet Kiosk"),
        ("Calendar & Scheduling", "Visual event timeline, date conflict detection, daily capacity load", "PC Back-Office"),
        ("Menu & Package Catalog", "CRUD packages, pax price matrix, dish categories, photo management", "PC + Tablet Kiosk"),
        ("Billing & Invoicing", "Invoice generation, downpayment enforcement, balance collection", "PC + Tablet Kiosk"),
        ("Receipt & Notifications", "PDF receipt generation, SMTP email dispatcher, SMS gateway", "PC + Tablet (PDF)"),
        ("Cash Flow Tracker", "Inflow/Outflow tracking, payment method breakdown, net balance", "PC Back-Office"),
        ("Expense & COGS Manager", "Food costs, labor wages, logistics, utilities, Net Profit margin", "PC Back-Office"),
        ("AI Demand & Forecaster", "Pax forecasting, ingredient quantity estimator, menu insights", "PC Back-Office"),
        ("Kitchen Kanban & Checklist", "Visual preparation cards, per-dish task checklist, order dispatch", "PC Back-Office"),
        ("Customer CRM & Loyalty", "Booking history, automated loyalty tiers (Bronze to VIP), follow-ups", "PC + Tablet Kiosk"),
        ("Reports & Analytics", "Revenue charts, package sales breakdown, PDF/Excel export", "PC Back-Office"),
        ("Security, Audit & Backups", "Full change audit logs, 1-click DB backup/restore, Excel data sync", "PC + Tablet Kiosk")
    ]
    for row_idx, (m, desc, plat) in enumerate(modules_data, start=1):
        c0 = table.cell(row_idx, 0)
        c1 = table.cell(row_idx, 1)
        c2 = table.cell(row_idx, 2)
        
        bg_col = COLOR_CARD_DARK if row_idx % 2 == 0 else RGBColor(20, 30, 48)
        for c in [c0, c1, c2]:
            c.fill.solid()
            c.fill.fore_color.rgb = bg_col

        p0 = c0.text_frame.paragraphs[0]
        p0.text = m
        p0.font.name = "Arial"
        p0.font.size = Pt(9.5)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_GOLD

        p1 = c1.text_frame.paragraphs[0]
        p1.text = desc
        p1.font.name = "Arial"
        p1.font.size = Pt(9)
        p1.font.color.rgb = COLOR_TEXT_LIGHT

        p2 = c2.text_frame.paragraphs[0]
        p2.text = plat
        p2.font.name = "Arial"
        p2.font.size = Pt(9)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 11: Technical Specifications & Security
    # ─────────────────────────────────────────────────────────────────────────
    s11 = add_base_slide("Technical Specifications & Architecture", "Engineering standards, offline reliability, and data privacy", "FULL SYSTEM")
    add_card(s11, Inches(0.8), Inches(1.7), Inches(5.6), Inches(4.9), "Technology Stack & Architecture", [
        "Back-Office PC System: Native Python/Qt Desktop Application compiled to standalone Windows .exe.",
        "Tablet Kiosk Application: HTML5 / Vanilla CSS / JS + Client-Side SQLite WASM + Standalone Android APK (7.1 MB).",
        "Offline-First Reliability: Embedded relational database engine requires zero internet connection.",
        "HD Media Processor: Full HD 1920x1440 bicubic compression for crisp food and package photos."
    ], accent_color=COLOR_GOLD)

    add_card(s11, Inches(6.9), Inches(1.7), Inches(5.6), Inches(4.9), "Security & Data Governance", [
        "Complete Change Audit Log: Records actor, action type, timestamp, and JSON data states.",
        "100% Data Privacy: Database stored locally on physical business hardware.",
        "Zero Cloud Vulnerabilities: Immune to external cloud server outages or data breaches.",
        "1-Click Backup & Sync: Automated SQLite .db backups and 2-way master Excel synchronization."
    ], accent_color=COLOR_CRIMSON)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 12: Conclusion & Q&A Defense
    # ─────────────────────────────────────────────────────────────────────────
    s12 = prs.slides.add_slide(blank_slide_layout)
    bg12 = s12.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg12.fill.solid()
    bg12.fill.fore_color.rgb = COLOR_BG_DARK
    bg12.line.fill.background()

    tb12 = s12.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.933), Inches(3.5))
    tf12 = tb12.text_frame
    tf12.word_wrap = True
    tf12.margin_left = tf12.margin_top = tf12.margin_right = tf12.margin_bottom = 0

    p = tf12.paragraphs[0]
    p.text = "THANK YOU!"
    p.font.name = "Arial"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_CRIMSON
    p.alignment = PP_ALIGN.CENTER

    p2 = tf12.add_paragraph()
    p2.text = "Jayraldine's Catering Management & Self-Service Kiosk System"
    p2.font.name = "Arial"
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_GOLD
    p2.space_before = Pt(14)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf12.add_paragraph()
    p3.text = "Floor is Open for Questions & Evaluation (Q&A Defense)"
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_TEXT_LIGHT
    p3.space_before = Pt(18)
    p3.alignment = PP_ALIGN.CENTER

    prs.save(PPTX_OUT)
    print(f"Generated PowerPoint Deck: {PPTX_OUT}")

if __name__ == "__main__":
    create_deck()
